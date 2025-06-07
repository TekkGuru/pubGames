#!/usr/bin/env python3
"""
PowerPoint Slide Recreation Script
Recreates the "Current Issues for Medical Sector" slide using python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import io
import base64

def create_medical_sector_slide():
    """
    Creates a PowerPoint slide matching the provided image
    """
    
    # Create presentation
    prs = Presentation()
    
    # Set slide dimensions (16:9 aspect ratio)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Add slide with blank layout
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Define colors
    teal_color = RGBColor(134, 171, 174)  # Teal/blue-green
    dark_text = RGBColor(64, 64, 64)      # Dark gray
    light_gray = RGBColor(245, 245, 245)  # Light gray for boxes
    
    # 1. ADD CURVED DESIGN ELEMENTS
    # Create curved shapes to match the design
    
    # Large curved element (bottom left)
    curve_shape1 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, 
        Inches(-2), Inches(4), 
        Inches(8), Inches(6)
    )
    curve_shape1.fill.solid()
    curve_shape1.fill.fore_color.rgb = teal_color
    curve_shape1.line.fill.background()
    
    # Medium curved element
    curve_shape2 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, 
        Inches(-1), Inches(2), 
        Inches(6), Inches(4)
    )
    curve_shape2.fill.solid()
    curve_shape2.fill.fore_color.rgb = RGBColor(160, 190, 193)
    curve_shape2.line.fill.background()
    
    # 2. ADD MAIN MEDICAL IMAGE
    # Note: In a real implementation, you would load the actual medical/surgery image
    # For this demo, we'll create a placeholder rectangle
    img_placeholder = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(0.5),
        Inches(6), Inches(4)
    )
    img_placeholder.fill.solid()
    img_placeholder.fill.fore_color.rgb = RGBColor(200, 220, 200)
    img_placeholder.line.color.rgb = RGBColor(150, 150, 150)
    
    # Add text to placeholder
    img_text = img_placeholder.text_frame
    img_text.text = "Medical Surgery Image\n(Operating Room Scene)"
    img_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    img_text.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # 3. ADD MAIN TITLE
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(5.2),
        Inches(6), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_frame.margin_left = 0
    title_frame.margin_right = 0
    
    # Title text
    title_p = title_frame.paragraphs[0]
    title_p.text = "Current Issues for Medical Sector"
    title_p.font.name = "Arial"
    title_p.font.size = Pt(36)
    title_p.font.bold = True
    title_p.font.color.rgb = dark_text
    title_p.alignment = PP_ALIGN.LEFT
    
    # 4. ADD SUBTITLE
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(6.5),
        Inches(6), Inches(0.8)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.text = "This slide is 100% editable. Adapt it to your needs and capture your audience's attention."
    subtitle_p.font.name = "Arial"
    subtitle_p.font.size = Pt(12)
    subtitle_p.font.color.rgb = RGBColor(100, 100, 100)
    subtitle_p.alignment = PP_ALIGN.LEFT
    
    # 5. ADD CONTENT BOXES
    
    # Box 1: Patients Changing Consumers Behavior
    box1 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(7.5), Inches(0.8),
        Inches(5), Inches(2)
    )
    box1.fill.solid()
    box1.fill.fore_color.rgb = light_gray
    box1.line.color.rgb = RGBColor(200, 200, 200)
    box1.line.width = Pt(1)
    
    box1_text = box1.text_frame
    box1_text.margin_left = Inches(0.2)
    box1_text.margin_right = Inches(0.2)
    box1_text.margin_top = Inches(0.1)
    box1_text.word_wrap = True
    
    # Box 1 Title
    p1_title = box1_text.paragraphs[0]
    p1_title.text = "Patients Changing Consumers Behavior"
    p1_title.font.name = "Arial"
    p1_title.font.size = Pt(14)
    p1_title.font.bold = True
    p1_title.font.color.rgb = dark_text
    
    # Box 1 Bullet 1
    p1_bullet1 = box1_text.add_paragraph()
    p1_bullet1.text = "Healthcare industry centers around convenience for patients and customer service."
    p1_bullet1.font.name = "Arial"
    p1_bullet1.font.size = Pt(10)
    p1_bullet1.font.italic = True
    p1_bullet1.font.color.rgb = dark_text
    p1_bullet1.level = 0
    
    # Box 1 Bullet 2
    p1_bullet2 = box1_text.add_paragraph()
    p1_bullet2.text = "Inculcate online health information and online request for appointments to provide more transparency to patients."
    p1_bullet2.font.name = "Arial"
    p1_bullet2.font.size = Pt(10)
    p1_bullet2.font.italic = True
    p1_bullet2.font.color.rgb = dark_text
    p1_bullet2.level = 0
    
    # Box 2: Marketing Limitations
    box2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(7.5), Inches(3),
        Inches(5), Inches(1.8)
    )
    box2.fill.solid()
    box2.fill.fore_color.rgb = light_gray
    box2.line.color.rgb = RGBColor(200, 200, 200)
    box2.line.width = Pt(1)
    
    box2_text = box2.text_frame
    box2_text.margin_left = Inches(0.2)
    box2_text.margin_right = Inches(0.2)
    box2_text.margin_top = Inches(0.1)
    box2_text.word_wrap = True
    
    # Box 2 Title
    p2_title = box2_text.paragraphs[0]
    p2_title.text = "Marketing Limitations"
    p2_title.font.name = "Arial"
    p2_title.font.size = Pt(14)
    p2_title.font.bold = True
    p2_title.font.color.rgb = dark_text
    
    # Box 2 Bullet 1
    p2_bullet1 = box2_text.add_paragraph()
    p2_bullet1.text = "As per government policies and ethical issues scope for advertising in the medical field is limited."
    p2_bullet1.font.name = "Arial"
    p2_bullet1.font.size = Pt(10)
    p2_bullet1.font.italic = True
    p2_bullet1.font.color.rgb = dark_text
    p2_bullet1.level = 0
    
    # Box 2 Bullet 2
    p2_bullet2 = box2_text.add_paragraph()
    p2_bullet2.text = "Healthcare leaders are utilizing pay per click advertising for the retargeting online website for users."
    p2_bullet2.font.name = "Arial"
    p2_bullet2.font.size = Pt(10)
    p2_bullet2.font.italic = True
    p2_bullet2.font.color.rgb = dark_text
    p2_bullet2.level = 0
    
    # Box 3: Optimal Staff Efficiency
    box3 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(7.5), Inches(5),
        Inches(5), Inches(2)
    )
    box3.fill.solid()
    box3.fill.fore_color.rgb = light_gray
    box3.line.color.rgb = RGBColor(200, 200, 200)
    box3.line.width = Pt(1)
    
    box3_text = box3.text_frame
    box3_text.margin_left = Inches(0.2)
    box3_text.margin_right = Inches(0.2)
    box3_text.margin_top = Inches(0.1)
    box3_text.word_wrap = True
    
    # Box 3 Title
    p3_title = box3_text.paragraphs[0]
    p3_title.text = "Optimal Staff Efficiency"
    p3_title.font.name = "Arial"
    p3_title.font.size = Pt(14)
    p3_title.font.bold = True
    p3_title.font.color.rgb = dark_text
    
    # Box 3 Bullet 1
    p3_bullet1 = box3_text.add_paragraph()
    p3_bullet1.text = "Due heavy patient load and low managerial experience skills of doctors cause havoc at hospital."
    p3_bullet1.font.name = "Arial"
    p3_bullet1.font.size = Pt(10)
    p3_bullet1.font.italic = True
    p3_bullet1.font.color.rgb = dark_text
    p3_bullet1.level = 0
    
    # Box 3 Bullet 2
    p3_bullet2 = box3_text.add_paragraph()
    p3_bullet2.text = "Seek or hire hospital management professionals just to manage and cope up with management activities on day to day basis."
    p3_bullet2.font.name = "Arial"
    p3_bullet2.font.size = Pt(10)
    p3_bullet2.font.italic = True
    p3_bullet2.font.color.rgb = dark_text
    p3_bullet2.level = 0
    
    return prs

def add_actual_image_to_slide(prs, image_path):
    """
    Function to add actual medical image if available
    
    Args:
        prs: PowerPoint presentation object
        image_path: Path to the medical/surgery image file
    """
    slide = prs.slides[0]
    
    # Remove placeholder and add actual image
    # Note: You would replace the placeholder rectangle with this
    try:
        slide.shapes.add_picture(
            image_path,
            Inches(0.5), Inches(0.5),
            width=Inches(6), height=Inches(4)
        )
        print(f"Added image: {image_path}")
    except Exception as e:
        print(f"Could not add image: {e}")

def enhance_slide_design(prs):
    """
    Additional design enhancements
    """
    slide = prs.slides[0]
    
    # Add more curved elements for better design matching
    # These would be fine-tuned based on the exact curvature needed
    
    # Additional decorative elements can be added here
    # Such as more precise curved shapes, gradients, etc.
    
    pass

def main():
    """
    Main function to create and save the PowerPoint presentation
    """
    print("Creating PowerPoint slide recreation...")
    
    # Create the slide
    presentation = create_medical_sector_slide()
    
    # Enhance design (optional)
    enhance_slide_design(presentation)
    
    # Save the presentation
    output_filename = "medical_sector_issues_slide.pptx"
    presentation.save(output_filename)
    
    print(f"✅ Slide created successfully: {output_filename}")
    print("\nSlide Elements Created:")
    print("- Curved background design elements")
    print("- Medical image placeholder")
    print("- Main title: 'Current Issues for Medical Sector'")
    print("- Subtitle with editable text")
    print("- Three content boxes with bullet points:")
    print("  1. Patients Changing Consumers Behavior")
    print("  2. Marketing Limitations") 
    print("  3. Optimal Staff Efficiency")
    
    print("\n📝 To customize further:")
    print("1. Replace image placeholder with actual medical/surgery image")
    print("2. Adjust colors to match exact brand colors")
    print("3. Fine-tune curved element positions")
    print("4. Modify text content as needed")

if __name__ == "__main__":
    main()


# ADDITIONAL UTILITY FUNCTIONS FOR IMAGE ANALYSIS

def analyze_image_elements(image_path):
    """
    Analyze the original image to extract element positions and properties
    This would use OpenCV, PIL, and OCR for detailed analysis
    """
    try:
        from PIL import Image, ImageDraw
        import cv2
        import numpy as np
        
        # Load image
        img = Image.open(image_path)
        img_cv = cv2.imread(image_path)
        
        # Convert to different color spaces for analysis
        hsv = cv2.cvtColor(img_cv, cv2.COLOR_BGR2HSV)
        gray = cv2.cvtColor(img_cv, cv2.COLOR_BGR2GRAY)
        
        analysis_results = {
            'image_dimensions': img.size,
            'detected_text_regions': [],
            'color_palette': [],
            'shapes_detected': [],
            'layout_analysis': {}
        }
        
        # Text region detection (would use OCR here)
        # This is a simplified version - real implementation would use Tesseract
        
        # Color analysis
        # Extract dominant colors
        
        # Shape detection
        # Detect rectangles, circles, etc.
        
        return analysis_results
        
    except ImportError:
        print("OpenCV and PIL required for detailed image analysis")
        return None

def extract_text_with_ocr(image_path):
    """
    Extract text from image using OCR
    """
    try:
        import pytesseract
        from PIL import Image
        
        img = Image.open(image_path)
        
        # Extract text with position information
        data = pytesseract.image_to_data(img, output_type=pytesseract.Output.DICT)
        
        text_elements = []
        for i in range(len(data['text'])):
            if int(data['conf'][i]) > 30:  # Confidence threshold
                text_elements.append({
                    'text': data['text'][i],
                    'x': data['left'][i],
                    'y': data['top'][i],
                    'width': data['width'][i],
                    'height': data['height'][i],
                    'confidence': data['conf'][i]
                })
        
        return text_elements
        
    except ImportError:
        print("pytesseract required for OCR functionality")
        return []

# ADVANCED FEATURES FOR EXACT RECREATION

def match_fonts_and_styles(text_elements):
    """
    Analyze and match fonts and text styles from the original image
    """
    # This would analyze text characteristics and map to available fonts
    font_mappings = {
        'title': {'name': 'Arial', 'size': 36, 'bold': True},
        'subtitle': {'name': 'Arial', 'size': 12, 'italic': False},
        'body': {'name': 'Arial', 'size': 10, 'italic': True},
        'headers': {'name': 'Arial', 'size': 14, 'bold': True}
    }
    return font_mappings

def create_precise_layout(analysis_data):
    """
    Create precise layout based on image analysis
    """
    # This would use the analysis data to position elements exactly
    # as they appear in the original image
    pass